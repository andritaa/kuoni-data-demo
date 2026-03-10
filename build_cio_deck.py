"""
Kuoni CIO Presentation — Stephen Adebola
Built with python-pptx
Brand: Teal #1B4F6B | Gold #C9A96E | White | Dark #0D1B24
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree
import copy

# ── Colours ──────────────────────────────────────────────
TEAL      = RGBColor(0x1B, 0x4F, 0x6B)
TEAL_LT   = RGBColor(0x2E, 0x6E, 0x8E)
GOLD      = RGBColor(0xC9, 0xA9, 0x6E)
GOLD_LT   = RGBColor(0xE0, 0xC8, 0x95)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
DARK      = RGBColor(0x0D, 0x1B, 0x24)
GREY      = RGBColor(0x6B, 0x72, 0x80)
LIGHT_BG  = RGBColor(0xF8, 0xF6, 0xF3)
LIGHT_GRY = RGBColor(0xE5, 0xE7, 0xEB)
MID_GRY   = RGBColor(0x9C, 0xA3, 0xAF)
GREEN     = RGBColor(0x16, 0xA3, 0x4A)
AMBER     = RGBColor(0xD9, 0x77, 0x06)
RED       = RGBColor(0xDC, 0x26, 0x26)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

blank_layout = prs.slide_layouts[6]   # completely blank

# ═══════════════════════════════════════════════════════════
# HELPERS
# ═══════════════════════════════════════════════════════════

def add_rect(slide, l, t, w, h, fill=None, line=None, line_w=Pt(0)):
    shape = slide.shapes.add_shape(1, l, t, w, h)   # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = line_w
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h,
             font_name='Calibri', font_size=Pt(12), bold=False, italic=False,
             color=WHITE, align=PP_ALIGN.LEFT, wrap=True, v_anchor=None):
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    if v_anchor:
        tf.auto_size = None
        from pptx.enum.text import MSO_ANCHOR
        tf.vertical_anchor = v_anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = font_size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox

def add_para(tf, text, font_name='Calibri', font_size=Pt(12), bold=False,
             italic=False, color=WHITE, align=PP_ALIGN.LEFT, space_before=Pt(0)):
    from pptx.oxml.ns import qn
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = space_before
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = font_size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return p

def add_textbox_multi(slide, lines, l, t, w, h, wrap=True):
    """lines = list of dicts with keys: text, size, bold, italic, color, align, space_before"""
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    first = True
    for ln in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = ln.get('align', PP_ALIGN.LEFT)
        if 'space_before' in ln:
            p.space_before = ln['space_before']
        run = p.add_run()
        run.text = ln.get('text', '')
        run.font.name = ln.get('font', 'Calibri')
        run.font.size = ln.get('size', Pt(12))
        run.font.bold = ln.get('bold', False)
        run.font.italic = ln.get('italic', False)
        run.font.color.rgb = ln.get('color', WHITE)
    return txBox

def add_logo(slide, x=Inches(11.8), y=Inches(0.18)):
    add_text(slide, 'KUONI', x, y, Inches(1.4), Inches(0.35),
             font_name='Georgia', font_size=Pt(16), bold=True, color=WHITE, align=PP_ALIGN.RIGHT)
    add_text(slide, 'Data Intelligence Platform', x, y+Inches(0.28), Inches(1.4), Inches(0.25),
             font_name='Calibri', font_size=Pt(8), color=GOLD, align=PP_ALIGN.RIGHT)

def header_bar(slide, label=''):
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.75), fill=TEAL)
    add_logo(slide)
    if label:
        add_text(slide, label.upper(), Inches(0.45), Inches(0.2), Inches(4), Inches(0.38),
                 font_size=Pt(9), bold=True, color=GOLD_LT, align=PP_ALIGN.LEFT)

def footer_bar(slide, text='Confidential · Stephen Adebola · March 2026'):
    add_rect(slide, 0, SLIDE_H - Inches(0.35), SLIDE_W, Inches(0.35), fill=TEAL)
    add_text(slide, text, Inches(0.4), SLIDE_H - Inches(0.3), Inches(9), Inches(0.28),
             font_size=Pt(7), color=RGBColor(0xA0, 0xBE, 0xD4), align=PP_ALIGN.LEFT)

def slide_title_block(slide, title, subtitle='', y=Inches(0.85)):
    add_textbox_multi(slide, [
        {'text': title, 'font': 'Georgia', 'size': Pt(30), 'bold': True,
         'color': TEAL, 'align': PP_ALIGN.LEFT},
    ], Inches(0.45), y, Inches(10), Inches(0.65))
    if subtitle:
        add_text(slide, subtitle, Inches(0.45), y + Inches(0.6), Inches(10), Inches(0.35),
                 font_size=Pt(11), color=GREY, align=PP_ALIGN.LEFT)
    # Gold underline
    add_rect(slide, Inches(0.45), y + Inches(1.1), Inches(1.2), Inches(0.04), fill=GOLD)

def card_box(slide, l, t, w, h, fill=LIGHT_BG, border=LIGHT_GRY, border_w=Pt(1)):
    r = add_rect(slide, l, t, w, h, fill=fill, line=border, line_w=border_w)
    return r

def card_accent_top(slide, l, t, w, color=TEAL):
    add_rect(slide, l, t, w, Inches(0.06), fill=color)

# ═══════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)

# Full teal background
add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=TEAL)

# Gold accent panel right
add_rect(slide, SLIDE_W - Inches(3.8), 0, Inches(3.8), SLIDE_H, fill=RGBColor(0x17, 0x44, 0x5D))

# Gold vertical bar
add_rect(slide, SLIDE_W - Inches(3.85), Inches(1.5), Inches(0.07), Inches(4.5), fill=GOLD)

# White horizontal line
add_rect(slide, Inches(0.6), Inches(3.6), Inches(8.2), Inches(0.03), fill=RGBColor(0xFF,0xFF,0xFF))

# Eyebrow
add_text(slide, 'CIO MEETING  ·  MARCH 2026', Inches(0.6), Inches(1.5), Inches(7.5), Inches(0.35),
         font_size=Pt(9), bold=True, color=GOLD, align=PP_ALIGN.LEFT)

# Main title
add_textbox_multi(slide, [
    {'text': 'Turning Snowflake', 'font': 'Georgia', 'size': Pt(42), 'bold': True, 'color': WHITE},
    {'text': 'into a Strategic Asset', 'font': 'Georgia', 'size': Pt(42), 'bold': True,
     'italic': True, 'color': GOLD},
], Inches(0.6), Inches(2.0), Inches(8.5), Inches(1.8))

# Subtitle
add_text(slide, 'A practical roadmap for making Kuoni\'s data platform the\nfoundation for better decisions and lasting competitive advantage.',
         Inches(0.6), Inches(3.75), Inches(7.8), Inches(0.9),
         font_size=Pt(13), color=RGBColor(0xA0,0xBE,0xD4), align=PP_ALIGN.LEFT)

# Meta box
add_rect(slide, Inches(0.6), Inches(5.0), Inches(5.5), Inches(1.8), fill=RGBColor(0x14,0x3D,0x57))
add_text(slide, 'Presented by', Inches(0.8), Inches(5.15), Inches(5), Inches(0.28),
         font_size=Pt(8), color=GOLD, bold=True)
add_text(slide, 'Stephen Adebola  —  Data Architect', Inches(0.8), Inches(5.4), Inches(5), Inches(0.3),
         font_size=Pt(13), color=WHITE, bold=True)
add_text(slide, 'For: Richard Nunn, CIO · Kuoni / DERTOUR Group', Inches(0.8), Inches(5.72), Inches(5), Inches(0.28),
         font_size=Pt(10), color=RGBColor(0xA0,0xBE,0xD4))
add_text(slide, 'Via: Xavier Labat · La Fosse Associates', Inches(0.8), Inches(5.98), Inches(5), Inches(0.28),
         font_size=Pt(10), color=RGBColor(0xA0,0xBE,0xD4))
add_text(slide, 'Strictly Confidential', Inches(0.8), Inches(6.35), Inches(5), Inches(0.28),
         font_size=Pt(8), color=GOLD)

# Right panel text
add_text(slide, 'KUONI', Inches(9.9), Inches(2.8), Inches(3), Inches(0.5),
         font_name='Georgia', font_size=Pt(28), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, 'Data Intelligence', Inches(9.9), Inches(3.3), Inches(3), Inches(0.35),
         font_size=Pt(12), color=GOLD, align=PP_ALIGN.CENTER)
add_text(slide, 'Platform', Inches(9.9), Inches(3.6), Inches(3), Inches(0.35),
         font_size=Pt(12), color=GOLD, align=PP_ALIGN.CENTER)
add_rect(slide, Inches(10.5), Inches(4.1), Inches(1.8), Inches(0.04), fill=GOLD)

# ═══════════════════════════════════════════════════════════
# SLIDE 2 — THE SITUATION
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=LIGHT_BG)
header_bar(slide, '01  ·  The Situation')
footer_bar(slide)
slide_title_block(slide, 'Where you are today', 'Snowflake is live — but working below its potential')

# LEFT CARD — What's Working
card_box(slide, Inches(0.45), Inches(2.0), Inches(3.9), Inches(4.6))
card_accent_top(slide, Inches(0.45), Inches(2.0), Inches(3.9), color=TEAL_LT)
add_text(slide, '✓  WHAT\'S WORKING', Inches(0.6), Inches(2.12), Inches(3.5), Inches(0.3),
         font_size=Pt(8), bold=True, color=TEAL)
items_working = [
    'Snowflake is deployed and running',
    'Initial reporting use case is functional',
    'Team has early Snowflake familiarity',
    'Foundation data is accessible',
    'Investment has been made and committed to',
]
for i, item in enumerate(items_working):
    add_text(slide, f'→  {item}', Inches(0.6), Inches(2.5) + Inches(0.47)*i, Inches(3.6), Inches(0.4),
             font_size=Pt(10.5), color=RGBColor(0x1F, 0x2D, 0x3D))

# MIDDLE CARD — The Gaps
card_box(slide, Inches(4.6), Inches(2.0), Inches(3.9), Inches(4.6), border=RED)
card_accent_top(slide, Inches(4.6), Inches(2.0), Inches(3.9), color=RED)
add_text(slide, '✗  THE GAPS', Inches(4.75), Inches(2.12), Inches(3.5), Inches(0.3),
         font_size=Pt(8), bold=True, color=RED)
items_gaps = [
    'No enterprise data architecture around it',
    'Data silos across booking, CRM & ops',
    'No single source of truth for customers',
    'No self-service for business teams',
    'No data governance or quality framework',
]
for i, item in enumerate(items_gaps):
    add_text(slide, f'→  {item}', Inches(4.75), Inches(2.5) + Inches(0.47)*i, Inches(3.6), Inches(0.4),
             font_size=Pt(10.5), color=RGBColor(0x1F, 0x2D, 0x3D))

# RIGHT CARD — The Opportunity
card_box(slide, Inches(8.75), Inches(2.0), Inches(4.15), Inches(4.6), fill=TEAL, border=TEAL)
add_text(slide, '→  THE OPPORTUNITY', Inches(8.9), Inches(2.12), Inches(3.8), Inches(0.3),
         font_size=Pt(8), bold=True, color=GOLD)
add_textbox_multi(slide, [
    {'text': 'Snowflake was a tactical fix.', 'font': 'Georgia', 'size': Pt(14), 'bold': True, 'color': WHITE},
    {'text': 'It can become your strategic advantage.', 'font': 'Georgia', 'size': Pt(14),
     'bold': True, 'italic': True, 'color': GOLD, 'space_before': Pt(4)},
], Inches(8.9), Inches(2.42), Inches(3.8), Inches(0.9))
items_opp = [
    'Unified view of every customer',
    'Real-time revenue & margin visibility',
    'Self-service analytics for all teams',
    'Foundation for personalisation at scale',
    'Data products for commercial decisions',
]
for i, item in enumerate(items_opp):
    add_text(slide, f'◆  {item}', Inches(8.9), Inches(3.55) + Inches(0.43)*i, Inches(3.8), Inches(0.38),
             font_size=Pt(10), color=RGBColor(0xC4, 0xD9, 0xE8))

add_text(slide, 'I\'ve done this before — at Marlink\nin an almost identical situation.',
         Inches(8.9), Inches(5.9), Inches(3.8), Inches(0.5),
         font_size=Pt(9), italic=True, color=RGBColor(0x80, 0xAA, 0xC4))

# ═══════════════════════════════════════════════════════════
# SLIDE 3 — THE OPPORTUNITY
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=LIGHT_BG)
header_bar(slide, '02  ·  The Opportunity')
footer_bar(slide)
slide_title_block(slide, 'What great looks like', 'What a properly designed data platform enables for a premium travel brand')

# 4 KPI boxes at top
kpis = [
    ('360°', 'Customer View', TEAL),
    ('Real-time', 'Revenue Visibility', RGBColor(0x0F, 0x76, 0x6E)),
    ('Self-service', 'Analytics for All Teams', RGBColor(0x92, 0x40, 0x0E)),
    ('Trusted', 'Single Source of Truth', RGBColor(0x1E, 0x3A, 0x5F)),
]
for i, (val, label, col) in enumerate(kpis):
    x = Inches(0.45) + Inches(3.1) * i
    card_box(slide, x, Inches(2.0), Inches(2.9), Inches(1.05), fill=col)
    add_text(slide, val, x, Inches(2.08), Inches(2.9), Inches(0.48),
             font_name='Georgia', font_size=Pt(22), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, label, x, Inches(2.54), Inches(2.9), Inches(0.35),
             font_size=Pt(9), color=GOLD_LT, align=PP_ALIGN.CENTER)

# 3 capability cards
caps = [
    ('✈️  Commercial Intelligence', TEAL, [
        'Booking patterns by destination, agent & season',
        'Revenue and margin per product line',
        'Cancellation and amendment trends',
        'Channel performance (online vs agent)',
        'Yield management and pricing insight',
    ]),
    ('👥  Customer Intelligence', RGBColor(0x78,0x35,0x06), [
        'Lifetime value and loyalty segmentation',
        'Repeat booking likelihood modelling',
        'Personalisation signals across touchpoints',
        'Churn risk identification before it happens',
        'Customer 360 — one view across all systems',
    ]),
    ('🏢  Operational Intelligence', RGBColor(0x14,0x53,0x2D), [
        'Supplier and destination performance',
        'Agent productivity and conversion rates',
        'Capacity utilisation and yield management',
        'DERTOUR Group consolidated reporting',
        'Real-time operational dashboards',
    ]),
]
for i, (title, col, items) in enumerate(caps):
    x = Inches(0.45) + Inches(4.22) * i
    card_box(slide, x, Inches(3.25), Inches(4.0), Inches(3.35))
    card_accent_top(slide, x, Inches(3.25), Inches(4.0), color=col)
    add_text(slide, title, x + Inches(0.15), Inches(3.38), Inches(3.7), Inches(0.35),
             font_size=Pt(10), bold=True, color=col)
    for j, item in enumerate(items):
        add_text(slide, f'→  {item}', x + Inches(0.15), Inches(3.82) + Inches(0.44)*j, Inches(3.7), Inches(0.38),
                 font_size=Pt(10), color=RGBColor(0x1F,0x2D,0x3D))

# ═══════════════════════════════════════════════════════════
# SLIDE 4 — ROADMAP
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=LIGHT_BG)
header_bar(slide, '03  ·  The Roadmap')
footer_bar(slide)
slide_title_block(slide, 'Getting there in three phases', 'Structured, low-risk progression — quick wins first, transformation second')

phases = [
    {
        'num': 'PHASE 1', 'title': 'Stabilise', 'timing': 'Months 1 – 2  ·  Foundations',
        'col': TEAL, 'light': RGBColor(0xEB, 0xF4, 0xF8),
        'items': [
            'Audit the existing Snowflake environment',
            'Map all data sources: booking, CRM, marketing, ops',
            'Design enterprise data model (Medallion architecture)',
            'Bronze → Silver → Gold layer implementation',
            'Data governance and naming conventions',
            'Connect Snowflake to Power BI — first trusted reports',
        ],
        'outcome': 'A clean, trusted foundation — no more data arguments in meetings',
    },
    {
        'num': 'PHASE 2', 'title': 'Elevate', 'timing': 'Months 3 – 4  ·  Value Creation',
        'col': RGBColor(0x92, 0x40, 0x0E), 'light': RGBColor(0xFD, 0xF6, 0xEC),
        'col': GOLD, 'light': RGBColor(0xFD, 0xF6, 0xEC),
        'items': [
            'Build Customer 360 — unified profile across all touchpoints',
            'Launch governed Data Products per business domain',
            'Automate data quality monitoring and alerting',
            'Enable self-service analytics via semantic layer',
            'Commercial revenue and margin dashboards live',
            'Integrate dbt for automated transformation pipelines',
        ],
        'outcome': 'Teams get answers in minutes, not weeks — without asking IT',
    },
    {
        'num': 'PHASE 3', 'title': 'Scale', 'timing': 'Months 5 – 6  ·  Strategic Advantage',
        'col': GREEN, 'light': RGBColor(0xF0, 0xFD, 0xF4),
        'items': [
            'DERTOUR Group integration — consolidated multi-brand view',
            'Predictive analytics: LTV scoring, churn risk, demand',
            'Real-time personalisation data feeds for CRM & marketing',
            'Snowflake credit governance and cost optimisation',
            'AI/ML-ready data products for next phase of innovation',
            'Centre of excellence — capability transferred to your team',
        ],
        'outcome': 'Kuoni is now a data-first business — built to compete and scale',
    },
]

for i, p in enumerate(phases):
    x = Inches(0.45) + Inches(4.22) * i
    # Card
    card_box(slide, x, Inches(2.1), Inches(4.0), Inches(4.7))
    card_accent_top(slide, x, Inches(2.1), Inches(4.0), color=p['col'])
    # Phase label
    add_text(slide, p['num'], x + Inches(0.15), Inches(2.22), Inches(3.7), Inches(0.28),
             font_size=Pt(8), bold=True, color=p['col'])
    # Title
    add_text(slide, p['title'], x + Inches(0.15), Inches(2.46), Inches(3.7), Inches(0.45),
             font_name='Georgia', font_size=Pt(20), bold=True, color=TEAL)
    # Timing
    add_text(slide, p['timing'], x + Inches(0.15), Inches(2.88), Inches(3.7), Inches(0.28),
             font_size=Pt(8.5), color=GREY)
    # Items
    for j, item in enumerate(p['items']):
        add_text(slide, f'◆  {item}', x + Inches(0.15), Inches(3.22) + Inches(0.39)*j, Inches(3.7), Inches(0.34),
                 font_size=Pt(9.5), color=RGBColor(0x1F,0x2D,0x3D))
    # Outcome box
    card_box(slide, x + Inches(0.15), Inches(5.6), Inches(3.7), Inches(0.9),
             fill=p['col'], border=p['col'])
    add_text(slide, f'📌  {p["outcome"]}',
             x + Inches(0.25), Inches(5.7), Inches(3.5), Inches(0.75),
             font_size=Pt(9), bold=False, color=WHITE)

# ═══════════════════════════════════════════════════════════
# SLIDE 5 — WHY ME
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=LIGHT_BG)
header_bar(slide, "04  ·  Why Me")
footer_bar(slide)
slide_title_block(slide, "I've done this before", "Directly relevant experience — mapped to Kuoni's exact situation")

# Experience cards LEFT
roles = [
    ('Marlink', 'Senior Data Architect  (Most Relevant)',
     'Led migration from Azure Databricks → Snowflake. Designed the enterprise data model for Customer, Product and Location domains. Architected MDM migration — exactly what Kuoni needs for its customer data. I\'ve been through this exact journey: tactical Snowflake → strategic platform.', TEAL),
    ('AllianzGI', 'Enterprise Platform Lead',
     'Administered the enterprise Databricks platform. Stepped into a complex, already-live environment and brought order quickly. Every pattern transfers to Snowflake: Medallion, dbt, governance, Unity Catalog ≈ Snowflake Horizon.', GOLD),
    ('NHS England / Home Office', 'Principal Data Engineer',
     '1,000+ pipeline migrations at NHS England. 500+ pipelines at the Home Office. Reported to senior leadership and delivered at government scale — comfortable translating architecture into board-level language.', GREEN),
]
for i, (co, role, detail, col) in enumerate(roles):
    y = Inches(2.1) + Inches(1.58) * i
    card_box(slide, Inches(0.45), y, Inches(7.0), Inches(1.45))
    add_rect(slide, Inches(0.45), y, Inches(0.07), Inches(1.45), fill=col)
    add_text(slide, co, Inches(0.65), y + Inches(0.1), Inches(6.6), Inches(0.32),
             font_size=Pt(13), bold=True, color=TEAL)
    add_text(slide, role, Inches(0.65), y + Inches(0.38), Inches(6.6), Inches(0.28),
             font_size=Pt(9), bold=True, color=col)
    add_text(slide, detail, Inches(0.65), y + Inches(0.62), Inches(6.6), Inches(0.75),
             font_size=Pt(9.5), color=RGBColor(0x37,0x40,0x51))

# Skills panel RIGHT
card_box(slide, Inches(7.75), Inches(2.1), Inches(5.1), Inches(4.7))
add_text(slide, 'CORE CAPABILITIES', Inches(7.95), Inches(2.22), Inches(4.7), Inches(0.28),
         font_size=Pt(8), bold=True, color=TEAL)

skills = [
    ('Snowflake Architecture & Design', 95),
    ('Data Modelling (ELDM / Star Schema)', 95),
    ('dbt / Medallion Architecture', 92),
    ('Terraform / Platform-as-Code', 90),
    ('Data Governance & Data Products', 88),
    ('Power BI / Analytics Layer', 83),
    ('C-suite Stakeholder Communication', 88),
]
for i, (skill, pct) in enumerate(skills):
    y = Inches(2.6) + Inches(0.54)*i
    add_text(slide, skill, Inches(7.95), y, Inches(3.6), Inches(0.28), font_size=Pt(9.5), color=RGBColor(0x1F,0x2D,0x3D))
    # track
    card_box(slide, Inches(7.95), y + Inches(0.28), Inches(4.7), Inches(0.12), fill=LIGHT_GRY, border=LIGHT_GRY)
    # fill
    fill_w = Inches(4.7) * pct / 100
    c = TEAL if pct > 90 else (GOLD if pct > 85 else TEAL_LT)
    add_rect(slide, Inches(7.95), y + Inches(0.28), fill_w, Inches(0.12), fill=c)

# Databricks callout
card_box(slide, Inches(7.75), Inches(6.42), Inches(5.1), Inches(0.7), fill=TEAL, border=TEAL)
add_text(slide, '⚡  Databricks expertise → transfers directly to Snowflake. Enterprise rigour, proven at scale.',
         Inches(7.9), Inches(6.48), Inches(4.8), Inches(0.55), font_size=Pt(9.5), color=WHITE)

# ═══════════════════════════════════════════════════════════
# SLIDE 6 — FIRST 90 DAYS
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=LIGHT_BG)
header_bar(slide, '05  ·  The Plan')
footer_bar(slide)
slide_title_block(slide, 'My first 90 days', 'A concrete, structured start — no long runways, no lengthy discovery phases')

periods = [
    ('Days 1–30', 'Listen & Assess', TEAL, [
        'Meet every team that touches data — what do they need, what frustrates them?',
        'Full audit of Snowflake: schemas, queries, credit consumption, pipelines',
        'Map all data sources and integration points',
        'Identify quick wins: what can we improve in 30 days?',
        'Understand DERTOUR Group reporting requirements',
        'Agree success metrics with Richard and the team',
    ], 'DELIVERABLE: Current State Report — honest assessment with prioritised recommendations'),
    ('Days 31–60', 'Build & Architect', GOLD, [
        'Implement Medallion architecture across existing data',
        'Deploy enterprise data model — Customer, Booking, Destination, Agent',
        'First trusted Power BI dashboards live (revenue, bookings, agents)',
        'Ship quick wins identified in Month 1 to stakeholders',
        'Data governance framework: ownership, naming, documentation',
        'Data quality checks automated — issues caught before reports break',
    ], 'DELIVERABLE: First trusted dashboards live — business teams can self-serve'),
    ('Days 61–90', 'Accelerate & Prove', GREEN, [
        'Data Products launched: Commercial, Customer, Operations domains',
        'dbt pipelines automating Silver and Gold layer transformations',
        'Data quality monitoring in place with stakeholder alerting',
        'Strategic roadmap for Phases 2 & 3 presented to leadership',
        'Business case for continued investment — ROI demonstrated with data',
        'Knowledge transfer — your team can maintain and extend the platform',
    ], 'DELIVERABLE: Strategic Data Roadmap — board-ready, with clear outcomes and timelines'),
]

for i, (badge, title, col, items, deliverable) in enumerate(periods):
    x = Inches(0.45) + Inches(4.22) * i
    card_box(slide, x, Inches(2.1), Inches(4.0), Inches(5.0))
    card_accent_top(slide, x, Inches(2.1), Inches(4.0), color=col)
    # Badge
    card_box(slide, x + Inches(0.15), Inches(2.22), Inches(1.5), Inches(0.28), fill=col, border=col)
    add_text(slide, badge, x + Inches(0.15), Inches(2.22), Inches(1.5), Inches(0.28),
             font_size=Pt(7.5), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, title, x + Inches(0.15), Inches(2.55), Inches(3.7), Inches(0.42),
             font_name='Georgia', font_size=Pt(17), bold=True, color=TEAL)
    for j, item in enumerate(items):
        add_text(slide, f'·  {item}', x + Inches(0.15), Inches(3.05) + Inches(0.38)*j, Inches(3.7), Inches(0.33),
                 font_size=Pt(9), color=RGBColor(0x1F,0x2D,0x3D))
    # Deliverable
    card_box(slide, x + Inches(0.12), Inches(5.38), Inches(3.76), Inches(0.85),
             fill=col, border=col)
    add_text(slide, deliverable, x + Inches(0.22), Inches(5.45), Inches(3.55), Inches(0.72),
             font_size=Pt(8.5), color=WHITE)

# ═══════════════════════════════════════════════════════════
# SLIDE 7 — ARCHITECTURE VISION
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=LIGHT_BG)
header_bar(slide, '06  ·  The Vision')
footer_bar(slide)
slide_title_block(slide, 'The target architecture', 'Where Snowflake sits in a modern, integrated data platform for Kuoni')

# LEFT — Architecture diagram
layers = [
    ('📥  Source Systems', 'Booking engine · CRM · Marketing · Finance · Agents', TEAL_LT),
    ('🥉  Bronze Layer', 'Raw ingestion — full history, immutable, schema-on-read', RGBColor(0xA1,0x62,0x07)),
    ('🥈  Silver Layer', 'Cleansed, standardised, deduplicated, validated', RGBColor(0x4B,0x55,0x63)),
    ('🥇  Gold Layer', 'Business-ready views · Data Products · KPIs · Aggregates', RGBColor(0x92,0x40,0x0E)),
    ('📊  Consumption', 'Power BI · Self-service Analytics · APIs · AI/ML ready', GREEN),
]
layer_h = Inches(0.72)
for i, (name, desc, col) in enumerate(layers):
    y = Inches(2.1) + (layer_h + Inches(0.08)) * i
    card_box(slide, Inches(0.45), y, Inches(6.0), layer_h, fill=col)
    add_text(slide, name, Inches(0.6), y + Inches(0.1), Inches(5.7), Inches(0.32),
             font_size=Pt(11), bold=True, color=WHITE)
    add_text(slide, desc, Inches(0.6), y + Inches(0.38), Inches(5.7), Inches(0.28),
             font_size=Pt(9), color=RGBColor(0xE5,0xEE,0xF5))
    if i < len(layers) - 1:
        add_text(slide, '↓', Inches(3.2), y + layer_h, Inches(0.3), Inches(0.1),
                 font_size=Pt(8), color=GREY, align=PP_ALIGN.CENTER)

# RIGHT — Principles
principles = [
    ('Governance at every layer', TEAL, [
        'Data ownership defined per domain',
        'Automated quality checks before Gold promotion',
        'Full lineage — know where every number comes from',
        'Role-based access — right data to right people',
    ]),
    ('Self-service for the business', GOLD, [
        'Marketing can answer their own questions',
        'Commercial builds their own reports',
        'IT focuses on architecture, not ad-hoc queries',
    ]),
    ('DERTOUR Group integration', GREEN, [
        'Kuoni\'s platform becomes the group model',
        'Consolidated reporting across all brands',
        'One architecture, scaled across the portfolio',
    ]),
]
y_offset = Inches(2.1)
for title, col, items in principles:
    h = Inches(0.32) + Inches(0.36) * len(items) + Inches(0.2)
    card_box(slide, Inches(6.75), y_offset, Inches(6.1), h)
    card_accent_top(slide, Inches(6.75), y_offset, Inches(6.1), color=col)
    add_text(slide, title, Inches(6.9), y_offset + Inches(0.1), Inches(5.8), Inches(0.28),
             font_size=Pt(9.5), bold=True, color=col)
    for j, item in enumerate(items):
        add_text(slide, f'◆  {item}', Inches(6.9), y_offset + Inches(0.38) + Inches(0.36)*j, Inches(5.8), Inches(0.3),
                 font_size=Pt(9.5), color=RGBColor(0x1F,0x2D,0x3D))
    y_offset += h + Inches(0.14)

# ═══════════════════════════════════════════════════════════
# SLIDE 8 — CLOSE / NEXT STEPS
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)

# Full teal bg
add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=TEAL)
add_rect(slide, SLIDE_W - Inches(4.0), 0, Inches(4.0), SLIDE_H, fill=RGBColor(0x17,0x44,0x5D))
add_rect(slide, SLIDE_W - Inches(4.05), Inches(1.5), Inches(0.07), Inches(4.5), fill=GOLD)

add_text(slide, 'NEXT STEPS', Inches(0.6), Inches(1.5), Inches(7.5), Inches(0.3),
         font_size=Pt(9), bold=True, color=GOLD)

add_textbox_multi(slide, [
    {'text': "Let's build something", 'font': 'Georgia', 'size': Pt(40), 'bold': True, 'color': WHITE},
    {'text': 'worth having.', 'font': 'Georgia', 'size': Pt(40), 'bold': True, 'italic': True, 'color': GOLD,
     'space_before': Pt(4)},
], Inches(0.6), Inches(2.0), Inches(8.5), Inches(1.8))

add_text(slide,
    'Kuoni has the data. Snowflake is in place. What\'s needed is the architecture,\n'
    'the model, and the expertise to connect them into something that actually\n'
    'drives the business forward.',
    Inches(0.6), Inches(3.85), Inches(7.8), Inches(1.0),
    font_size=Pt(12), color=RGBColor(0xA0,0xBE,0xD4))

# 3 proposal boxes
props = [
    ('ENGAGEMENT', 'Data Architect\n3 – 6 months'),
    ('FIRST STEP', '2-week discovery\nCurrent state audit'),
    ('90-DAY OUTCOME', 'Strategic roadmap\nboard-ready'),
]
for i, (label, val) in enumerate(props):
    x = Inches(0.6) + Inches(2.4) * i
    card_box(slide, x, Inches(5.1), Inches(2.2), Inches(1.2), fill=RGBColor(0x14,0x3D,0x57), border=RGBColor(0x2E,0x6E,0x8E))
    add_text(slide, label, x, Inches(5.18), Inches(2.2), Inches(0.25),
             font_size=Pt(7.5), bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    add_text(slide, val, x, Inches(5.44), Inches(2.2), Inches(0.6),
             font_size=Pt(11), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Contact
add_text(slide, 'Stephen Adebola   ·   stephen@haba.io   ·   Agent: Xavier Labat, La Fosse',
         Inches(0.6), Inches(6.5), Inches(8), Inches(0.35),
         font_size=Pt(9), color=RGBColor(0x80,0xAA,0xC4))

# Right logo
add_text(slide, 'KUONI', Inches(9.9), Inches(2.8), Inches(3), Inches(0.5),
         font_name='Georgia', font_size=Pt(28), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, 'Data Intelligence', Inches(9.9), Inches(3.32), Inches(3), Inches(0.3),
         font_size=Pt(12), color=GOLD, align=PP_ALIGN.CENTER)
add_text(slide, 'Platform', Inches(9.9), Inches(3.6), Inches(3), Inches(0.3),
         font_size=Pt(12), color=GOLD, align=PP_ALIGN.CENTER)
add_rect(slide, Inches(10.5), Inches(4.1), Inches(1.8), Inches(0.04), fill=GOLD)


# ── SAVE ──────────────────────────────────────────────────
out = '/home/stephen/projects/kuoni-data-demo/Kuoni_CIO_Presentation_Stephen_Adebola.pptx'
prs.save(out)
print(f'Saved: {out}')
